import copy
import io
import json
import os
import re
import time
import base64
import requests  # <-- Added for AI Pipe support
from typing import Any, Dict, List, Optional, Tuple

from lxml import etree as _etree

from fastapi import FastAPI, UploadFile, Form, HTTPException
from fastapi.responses import HTMLResponse, StreamingResponse, FileResponse, Response

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------- Optional SDKs ----------
try:
    from openai import OpenAI  # OpenAI-style client (also works for AI Pipe via base_url)
except Exception:
    OpenAI = None

try:
    import anthropic  # Anthropic (Claude)
except Exception:
    anthropic = None

try:
    from google import genai  # Google Gemini
except Exception:
    genai = None

# ---------- App ----------
app = FastAPI(title="Your Text, Your Style – PPTX Generator")

# ---------- Limits / Defaults ----------
MAX_TEXT_CHARS = 60_000
MIN_SLIDES = 10
MAX_SLIDES = 40
MAX_TEMPLATE_BYTES = 30 * 1024 * 1024  # 30 MB
OPENAI_DEFAULT_MODEL = "gpt-4o-mini"        # good default for OpenAI / AI Pipe
ANTHROPIC_DEFAULT_MODEL = "claude-3-5-sonnet-latest"
GEMINI_DEFAULT_MODEL = "gemini-2.5-flash"

# ---------- Front page ----------
@app.get("/", response_class=HTMLResponse)
async def serve_frontend():
    """Serve the main HTML interface"""
    html_path = os.path.join(os.path.dirname(__file__), "index.html")
    try:
        with open(html_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    except FileNotFoundError:
        return HTMLResponse(
            content="<h1>Frontend not found</h1><p>Please ensure index.html is next to app.py</p>",
            status_code=404,
        )

# ---------- Favicon (with fallback) ----------
_FAVICON_FALLBACK_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO3n+9QAAAAASUVORK5CYII="
)

@app.get("/favicon.ico", include_in_schema=False)
async def favicon():
    """
    Serve favicon.ico if present in the working directory.
    Otherwise return a tiny transparent PNG to avoid 404s.
    """
    path = "favicon.ico"
    if os.path.exists(path):
        return FileResponse(path, media_type="image/x-icon")
    return Response(content=_FAVICON_FALLBACK_PNG, media_type="image/png")

# ---------- Generate endpoint ----------
@app.post("/generate")
async def generate_pptx(
    text: str = Form(...),
    guidance: Optional[str] = Form(None),
    provider: str = Form(...),                # "openai", "aipipe", "anthropic", "gemini"
    api_key: str = Form(...),                 # OpenAI key / AI Pipe token / Gemini key / Anthropic key
    model: Optional[str] = Form(None),
    num_slides: Optional[int] = Form(None),   # <-- NEW: desired number of slides (1..40)
    reuse_images: bool = Form(False),         # copy exact images from uploaded PPT slides
    template: Optional[UploadFile] = None,    # OPTIONAL template (.pptx/.potx)
):
    # Validate text
    if not text or not text.strip():
        raise HTTPException(status_code=400, detail="Text is required.")

    # Normalize & clamp target slides
    target_slides = None
    if num_slides is not None:
        try:
            target_slides = max(1, min(MAX_SLIDES, int(num_slides)))
        except Exception:
            target_slides = MIN_SLIDES

    # Optional template handling (validate if provided)
    tpl_bytes: Optional[bytes] = None
    if template and template.filename:
        if not template.filename.lower().endswith((".pptx", ".potx")):
            raise HTTPException(status_code=400, detail="Template must be .pptx or .potx.")
        tpl_bytes = await template.read()
        if len(tpl_bytes) < 1024:
            raise HTTPException(status_code=400, detail="Template looks empty or invalid.")
        if len(tpl_bytes) > MAX_TEMPLATE_BYTES:
            raise HTTPException(status_code=400, detail="Template too large (max 30 MB).")

    # Build slide plan using chosen provider (JSON-only)
    try:
        plan = await build_slide_plan_with_retry(
            text=text.strip()[:MAX_TEXT_CHARS],
            guidance=(guidance or "").strip(),
            provider=provider.lower().strip(),
            api_key=api_key.strip(),
            model=(model or "").strip() or None,
            target_slides=target_slides,   # pass through to instruct LLM
            max_retries=2,
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM error: {e}")

    # Enforce slide count (exact if provided; otherwise ensure minimum)
    if target_slides:
        plan = enforce_target_slides(plan, target=target_slides, max_slides=MAX_SLIDES)
    else:
        plan = ensure_min_slides(plan, min_slides=MIN_SLIDES, max_slides=MAX_SLIDES)

    # Fill any slides that still have no bullets with a targeted LLM call
    try:
        plan = await _fill_missing_bullets(
            plan=plan,
            context=text.strip()[:MAX_TEXT_CHARS],
            provider=provider.lower().strip(),
            api_key=api_key.strip(),
            model=(model or "").strip() or None,
        )
    except Exception:
        pass  # never block the response

    # Build PPTX from plan + template style
    try:
        out_bytes = build_presentation_from_plan(
            template_bytes=tpl_bytes,
            plan=plan,
            exact_reuse_images=bool(reuse_images),
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PowerPoint build error: {e}")

    headers = {
        "Content-Disposition": 'attachment; filename="generated.pptx"',
        "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "Cache-Control": "no-store",
    }
    return StreamingResponse(io.BytesIO(out_bytes), headers=headers)

# ---------- Slide Plan (LLM) ----------
SLIDE_SCHEMA_EXAMPLE = {
    "slides": [
        {"title": "A short slide title", "bullets": ["point 1", "point 2", "point 3"]}
    ]
}

def _llm_instruction(text: str, guidance: str, target_slides: Optional[int]) -> str:
    count_line = (
        f"- Choose exactly {target_slides} slides (if content is short, expand; if long, summarize).\n"
        if target_slides else
        f"- Choose a reasonable number of slides (min {MIN_SLIDES}, max {MAX_SLIDES})\n"
    )
    example = json.dumps({
        "slides": [
            {"title": "What It Is", "bullets": ["Specific fact from the text", "Another key detail", "Third insight"]},
            {"title": "How It Works", "bullets": ["Step one description", "Step two description", "Step three description"]},
        ]
    }, indent=2)
    return f"""You are a PowerPoint slide planner. Your ONLY output is a valid JSON object — no prose, no markdown, no code fences.

STRICT RULES (violating any rule makes the output unusable):
{count_line}
1. The JSON must have exactly this shape: {{"slides": [{{"title": "...", "bullets": ["...","...","..."]}}]}}
2. EVERY slide MUST have a "bullets" array with AT LEAST 3 and AT MOST 6 items.
3. A slide with an empty "bullets" array or fewer than 3 bullets is INVALID. Do NOT produce such slides.
4. Each bullet must be a complete, informative sentence or phrase drawn from the TEXT below. No placeholders, no "TBD", no "point 1".
5. Title ≤ 80 chars. Each bullet ≤ 120 chars.
6. If the TEXT is short, expand each idea into multiple slides with detailed bullets.
7. If a slide topic has no obvious content, infer reasonable supporting points from context.

Tone / style guidance: "{guidance}"

EXAMPLE of correct output:
{example}

TEXT TO CONVERT INTO SLIDES:
{text}""".strip()

async def build_slide_plan_with_retry(
    text: str,
    guidance: str,
    provider: str,
    api_key: str,
    model: Optional[str],
    target_slides: Optional[int],
    max_retries: int = 2,
) -> Dict[str, Any]:
    last_err: Optional[Exception] = None
    for attempt in range(max_retries + 1):
        try:
            return await build_slide_plan(text, guidance, provider, api_key, model, target_slides)
        except Exception as e:
            last_err = e
            if attempt == max_retries:
                break
            time.sleep(0.8 * (attempt + 1))  # small backoff
    assert last_err is not None
    raise last_err

async def _fill_missing_bullets(
    plan: Dict[str, Any],
    context: str,
    provider: str,
    api_key: str,
    model: Optional[str],
) -> Dict[str, Any]:
    """
    For any slide in `plan` whose bullets list is empty, call the LLM again
    with a targeted prompt to generate 3–5 bullets, then patch them back in.
    """
    slides = plan.get("slides", [])
    empty_idxs = [
        i for i, s in enumerate(slides)
        if not [b for b in (s.get("bullets") or []) if str(b).strip()]
    ]
    if not empty_idxs:
        return plan

    titles_block = "\n".join(
        f'{j+1}. {slides[i]["title"]}' for j, i in enumerate(empty_idxs)
    )
    fill_text = (
        f"CONTEXT (use this to write the bullets):\n{context[:3000]}\n\n"
        f"The {len(empty_idxs)} slide title(s) below each need exactly 3–5 bullets. "
        f"Write real, informative bullet points from the CONTEXT above. "
        f"Return ONLY a JSON object with a 'slides' array — no empty bullets allowed:\n\n"
        f"{titles_block}"
    )
    try:
        filled_plan = await build_slide_plan(
            text=fill_text,
            guidance="Every slide MUST have 3-5 bullets. Use the context to write specific, informative content.",
            provider=provider,
            api_key=api_key,
            model=model,
            target_slides=len(empty_idxs),
        )
        filled_slides = filled_plan.get("slides", [])
        for j, slide_idx in enumerate(empty_idxs):
            if j < len(filled_slides):
                new_bullets = [
                    str(b).strip()
                    for b in (filled_slides[j].get("bullets") or [])
                    if str(b).strip()
                ]
                if new_bullets:
                    slides[slide_idx]["bullets"] = new_bullets
    except Exception:
        pass  # never let fill-in break the whole request
    return plan

# ---------- AI Pipe call function ----------
def call_aipipe(api_key, model_name, messages, max_tokens=4096):
    """
    Calls AI Pipe's OpenRouter-compatible API for generating completions.
    """
    url = "https://aipipe.org/openrouter/v1/chat/completions"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key.strip()}",
    }

    body = {
        "model": model_name,
        "messages": messages,
        "max_tokens": max_tokens,
        "temperature": 0.7
    }

    response = requests.post(url, headers=headers, json=body)
    if response.status_code != 200:
        raise RuntimeError(f"AI Pipe error {response.status_code}: {response.text}")

    data = response.json()

    # Safely extract message content
    try:
        content = data["choices"][0]["message"]["content"]
        slides = json.loads(content)
        if "slides" not in slides:
            raise ValueError("No slides array in AI Pipe response")
        return slides
    except Exception:
        # Fallback: If AI Pipe sends plain text instead of JSON
        return {"slides": [{"title": "Generated Deck", "bullets": [content]}]}

async def build_slide_plan(
    text: str,
    guidance: str,
    provider: str,
    api_key: str,
    model: Optional[str],
    target_slides: Optional[int],
) -> Dict[str, Any]:
    model_name = (
        model
        or (OPENAI_DEFAULT_MODEL if provider in ("openai", "aipipe") else
            ANTHROPIC_DEFAULT_MODEL if provider == "anthropic" else
            GEMINI_DEFAULT_MODEL)
    )

    instruction = _llm_instruction(text, guidance, target_slides)
    messages = [{"role": "user", "content": instruction}]

    if provider == "aipipe":
        # Use direct AI Pipe call, do NOT use response_format
        return call_aipipe(api_key, model_name, messages)

    if provider == "openai":
        if OpenAI is None:
            raise RuntimeError("openai package not installed. `pip install openai`")
        client_kwargs = {"api_key": api_key}
        client = OpenAI(**client_kwargs)
        resp = client.responses.create(
            model=model_name,
            input=[{"role": "user", "content": instruction}],

            temperature=0.2,
        )
        content = _extract_openai_output_text(resp)
        data = _safe_json_parse(content)

    elif provider == "anthropic":
        if anthropic is None:
            raise RuntimeError("anthropic package not installed. `pip install anthropic`")
        client = anthropic.Anthropic(api_key=api_key)
        msg = client.messages.create(
            model=model_name,
            max_tokens=4096,
            system="Return ONLY valid JSON. No explanations. No code fences.",
            messages=messages,
            temperature=0.2,
        )
        text_out = "".join(
            blk.text for blk in getattr(msg, "content", []) if getattr(blk, "type", "") == "text"
        )
        data = _safe_json_parse(text_out)

    elif provider == "gemini":
        if genai is None:
            raise RuntimeError("google-genai package not installed. `pip install google-genai`")
        client = genai.Client(api_key=api_key)
        resp = client.models.generate_content(
            model=model_name,
            contents=instruction,
            config={"response_mime_type": "application/json"},
        )
        data = _safe_json_parse(resp.text or "")

    else:
        raise HTTPException(status_code=400, detail=f"Unsupported provider: {provider}")

    if not isinstance(data, dict) or "slides" not in data or not isinstance(data["slides"], list):
        raise RuntimeError("Provider did not return JSON with a 'slides' array.")
    return data

def _extract_openai_output_text(resp: Any) -> str:
    txt = getattr(resp, "output_text", None)
    if txt:
        return txt
    try:
        out = getattr(resp, "output", None)
        if out and len(out) and hasattr(out[0], "content") and len(out[0].content):
            maybe = out[0].content[0]
            if hasattr(maybe, "text"):
                return maybe.text
    except Exception:
        pass
    try:
        rs = getattr(resp, "responses", None)
        if rs and len(rs) and hasattr(rs[0], "output_text"):
            return rs[0].output_text
    except Exception:
        pass
    return json.dumps(resp, default=str)

def _safe_json_parse(s: str) -> Dict[str, Any]:
    s = (s or "").strip()
    if not s:
        return {"slides": []}
    try:
        return json.loads(s)
    except Exception:
        m = re.search(r"\{.*\}", s, flags=re.S)
        if m:
            return json.loads(m.group(0))
        raise

# ---------- Enforce slide counts ----------
def ensure_min_slides(plan: Dict[str, Any], min_slides: int, max_slides: int) -> Dict[str, Any]:
    """Ensure the plan has at least `min_slides` (<= max_slides)."""
    slides = plan.get("slides") or []
    out: List[Dict[str, Any]] = []
    for s in slides:
        title = str(s.get("title", "")).strip() or "Slide"
        bullets = [str(b).strip() for b in (s.get("bullets") or []) if str(b).strip()]
        out.append({"title": title, "bullets": bullets})

    # Split dense slides into chunks of up to 3 bullets
    i = 0
    while len(out) < min_slides and i < len(out) and len(out) < max_slides:
        s = out[i]
        if len(s["bullets"]) > 3:
            extra = s["bullets"][3:]
            s["bullets"] = s["bullets"][:3]
            while extra and len(out) < min_slides and len(out) < max_slides:
                chunk = extra[:3]
                extra = extra[3:]
                out.insert(i + 1, {"title": f"{s['title']} (cont.)", "bullets": chunk})
                i += 1
        i += 1

    # Pad with title-only slides if still fewer than min
    while len(out) < min_slides and len(out) < max_slides:
        out.append({"title": f"Slide {len(out)+1}", "bullets": []})

    plan["slides"] = out[:max_slides]
    return plan

def enforce_target_slides(plan: Dict[str, Any], target: int, max_slides: int) -> Dict[str, Any]:
    """Force the plan to have exactly `target` slides (clamped to max_slides)."""
    target = max(1, min(max_slides, int(target)))
    slides = plan.get("slides") or []
    # Normalize
    norm: List[Dict[str, Any]] = []
    for s in slides:
        norm.append({
            "title": (str(s.get("title", "")) or "Slide").strip(),
            "bullets": [str(b).strip() for b in (s.get("bullets") or []) if str(b).strip()],
        })

    # If too few: split & pad
    if len(norm) < target:
        norm_wrap = ensure_min_slides({"slides": norm}, min_slides=target, max_slides=max_slides)["slides"]
        plan["slides"] = norm_wrap[:target]
        return plan

    # If too many: try to merge simple "(cont.)" slides, else truncate
    if len(norm) > target:
        merged: List[Dict[str, Any]] = []
        i = 0
        while i < len(norm):
            cur = norm[i]
            if i + 1 < len(norm) and norm[i + 1]["title"].startswith(cur["title"]):
                # merge a continuation into current (cap bullets ~8)
                nxt = norm[i + 1]
                cur["bullets"].extend(nxt["bullets"])
                cur["bullets"] = cur["bullets"][:8]
                i += 2
                merged.append(cur)
            else:
                merged.append(cur)
                i += 1
        norm = merged

    plan["slides"] = norm[:target]
    return plan

# ---------- Slide background & decoration helpers ----------
_PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
_DML = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

def _get_slide_bg_elem(slide):
    """Return a deep-copy of the <p:bg> element from a slide, or None."""
    cSld = slide._element.find(f'{{{_PML}}}cSld')
    if cSld is not None:
        bg = cSld.find(f'{{{_PML}}}bg')
        if bg is not None:
            return copy.deepcopy(bg)
    return None

def _apply_slide_bg(slide, bg_elem):
    """Insert (or replace) a <p:bg> element into a slide's cSld, before spTree."""
    cSld = slide._element.find(f'{{{_PML}}}cSld')
    if cSld is None:
        return
    existing = cSld.find(f'{{{_PML}}}bg')
    if existing is not None:
        cSld.remove(existing)
    spTree = cSld.find(f'{{{_PML}}}spTree')
    pos = list(cSld).index(spTree) if spTree is not None else 0
    cSld.insert(pos, copy.deepcopy(bg_elem))

def _has_r_ref(elem) -> bool:
    """True if element or any descendant has an r:embed/r:id/r:link pointing to a rId."""
    for sub in elem.iter():
        for attr in sub.attrib:
            local = attr.split('}')[-1] if '}' in attr else attr
            if local in ('embed', 'id', 'link', 'href') and sub.attrib[attr].startswith('rId'):
                return True
    return False

def _sp_has_text(sp_elem) -> bool:
    """True if a <p:sp> contains any non-whitespace text."""
    for t in sp_elem.findall(f'.//{{{_DML}}}t'):
        if t.text and t.text.strip():
            return True
    return False

def _elem_has_any_text(elem) -> bool:
    """True if an element (sp, grpSp, graphicFrame, …) has any non-whitespace text anywhere."""
    for t in elem.findall(f'.//{{{_DML}}}t'):
        if t.text and t.text.strip():
            return True
    return False

def _is_slide_number_shape(sp_elem) -> bool:
    """True if the shape looks like a page/slide number (short numeric string, tiny area)."""
    texts = [t.text for t in sp_elem.findall(f'.//{{{_DML}}}t')
             if t.text and t.text.strip()]
    full = ''.join(texts).strip()
    if not full:
        return False
    # Pure number ≤ 3 chars (slide numbers like "1", "12"), or very tiny shapes
    if full.isdigit() and len(full) <= 3:
        return True
    xfrm = _get_sp_xfrm(sp_elem)
    if xfrm and xfrm[2] * xfrm[3] < 200_000_000:   # < ~0.015 sq-inch in EMU²
        return True
    return False

def _get_sp_xfrm(sp_elem) -> Optional[Tuple[int, int, int, int]]:
    """Return (left, top, width, height) in EMU for a <p:sp>, or None."""
    spPr = sp_elem.find(f'{{{_PML}}}spPr')
    if spPr is None:
        return None
    xfrm = spPr.find(f'{{{_DML}}}xfrm')
    if xfrm is None:
        return None
    off = xfrm.find(f'{{{_DML}}}off')
    ext = xfrm.find(f'{{{_DML}}}ext')
    if off is None or ext is None:
        return None
    try:
        return (int(off.get('x', 0)), int(off.get('y', 0)),
                int(ext.get('cx', 0)), int(ext.get('cy', 0)))
    except Exception:
        return None

def _get_decorative_elems(slide) -> List[Any]:
    """
    Return deep-copies of purely visual shapes (no text, no external refs).
    - Skips spTree metadata (nvGrpSpPr, grpSpPr)
    - Skips any <p:sp> with text  (those become title/body)
    - Skips any <p:grpSp> that contains text  (content grids, icon cards with labels)
    - Skips anything with external relationship refs (images, charts)
    Only pure vector/freeform shapes and text-free groups survive.
    """
    result: List[Any] = []
    cSld = slide._element.find(f'{{{_PML}}}cSld')
    if cSld is None:
        return result
    spTree = cSld.find(f'{{{_PML}}}spTree')
    if spTree is None:
        return result
    for elem in spTree:
        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
        if tag in ('nvGrpSpPr', 'grpSpPr'):
            continue
        if _has_r_ref(elem):
            continue                        # images / charts
        if tag == 'sp' and _sp_has_text(elem):
            continue                        # text boxes → handled as title/body
        if tag == 'grpSp' and _elem_has_any_text(elem):
            continue                        # content groups (icon cards, data grids, etc.)
        result.append(copy.deepcopy(elem))
    return result

def _insert_decorative_elems(slide, elems: List[Any]) -> None:
    """Insert decorative elements at the bottom of z-order (behind text)."""
    cSld = slide._element.find(f'{{{_PML}}}cSld')
    if cSld is None:
        return
    spTree = cSld.find(f'{{{_PML}}}spTree')
    if spTree is None:
        return
    insert_pos = 0
    for i, child in enumerate(spTree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('nvGrpSpPr', 'grpSpPr'):
            insert_pos = i + 1
        else:
            break
    for i, elem in enumerate(elems):
        spTree.insert(insert_pos + i, copy.deepcopy(elem))

def _get_template_text_shapes(slide) -> Tuple[Optional[Any], Optional[Any]]:
    """
    Identify the title and body text shapes in a template slide.

    Strategy (in priority order):
    1. Shapes with <p:ph type="title"> or <p:ph type="ctrTitle"> → title
    2. Shapes with <p:ph type="body"> or <p:ph type="subTitle"> → body
    3. Non-placeholder freeform text boxes (Canva style):
       sort by top-position; topmost = title, next = body
    Skips slide/page numbers and tiny shapes throughout.
    """
    _TITLE_PH_TYPES = {'title', 'ctrTitle'}
    _BODY_PH_TYPES  = {'body', 'subTitle'}

    cSld = slide._element.find(f'{{{_PML}}}cSld')
    if cSld is None:
        return None, None
    spTree = cSld.find(f'{{{_PML}}}spTree')
    if spTree is None:
        return None, None

    title_sp: Optional[Any] = None
    body_sp:  Optional[Any] = None
    freeform_candidates: List[Tuple[int, int, Any]] = []   # (top_y, area, elem)

    for elem in spTree:
        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
        if tag != 'sp':
            continue
        if not _sp_has_text(elem):
            continue
        if _is_slide_number_shape(elem):
            continue
        xfrm = _get_sp_xfrm(elem)
        if xfrm is None:
            continue
        area = xfrm[2] * xfrm[3]
        if area < 300_000_000:              # skip tiny labels
            continue

        ph = elem.find(f'.//{{{_PML}}}ph')
        if ph is not None:
            ph_type = ph.get('type') or ''  # '' means default (body) in OOXML
            if ph_type in _TITLE_PH_TYPES:
                if title_sp is None:
                    title_sp = elem
            elif ph_type in _BODY_PH_TYPES or ph_type == '':
                if body_sp is None:
                    body_sp = elem
        else:
            # Canva-style freeform text box — collect for position-based fallback
            freeform_candidates.append((xfrm[1], area, elem))  # top_y, area, elem

    # Fill missing slots from freeform candidates.
    # For Canva-style shapes (no placeholder type), sort by area DESC:
    # the largest text shape is the main headline, second-largest is body.
    # This prevents small branding/logo text boxes from being chosen as titles.
    if freeform_candidates:
        freeform_candidates.sort(key=lambda x: -x[1])   # area descending
        for _, _, elem in freeform_candidates:
            if title_sp is None:
                title_sp = elem
            elif body_sp is None:
                body_sp = elem
            if title_sp is not None and body_sp is not None:
                break

    return title_sp, body_sp

def _clone_sp_with_text(sp_elem, text: str) -> Any:
    """
    Deep-copy a Canva <p:sp>, clear its paragraphs, set a single paragraph
    with `text` while preserving the first run's formatting (font/size/color).
    """
    elem = copy.deepcopy(sp_elem)
    txBody = elem.find(f'{{{_PML}}}txBody')
    if txBody is None:
        return elem
    paras = txBody.findall(f'{{{_DML}}}p')
    first_pPr = None
    first_rPr = None
    for para in paras:
        pPr = para.find(f'{{{_DML}}}pPr')
        if pPr is not None and first_pPr is None:
            first_pPr = copy.deepcopy(pPr)
        for run in para.findall(f'{{{_DML}}}r'):
            rPr = run.find(f'{{{_DML}}}rPr')
            if rPr is not None and first_rPr is None:
                first_rPr = copy.deepcopy(rPr)
    for p in list(paras):
        txBody.remove(p)
    new_p = _etree.SubElement(txBody, f'{{{_DML}}}p')
    if first_pPr is not None:
        new_p.append(first_pPr)
    new_r = _etree.SubElement(new_p, f'{{{_DML}}}r')
    if first_rPr is not None:
        new_r.append(first_rPr)
    new_t = _etree.SubElement(new_r, f'{{{_DML}}}t')
    new_t.text = text
    return elem

def _clone_sp_with_bullets(sp_elem, bullets: List[str]) -> Any:
    """
    Deep-copy a Canva <p:sp>, replace its text with one paragraph per bullet,
    preserving the first run's formatting for all bullets.
    """
    elem = copy.deepcopy(sp_elem)
    txBody = elem.find(f'{{{_PML}}}txBody')
    if txBody is None:
        return elem
    paras = txBody.findall(f'{{{_DML}}}p')
    first_pPr = None
    first_rPr = None
    for para in paras:
        pPr = para.find(f'{{{_DML}}}pPr')
        if pPr is not None and first_pPr is None:
            first_pPr = copy.deepcopy(pPr)
        for run in para.findall(f'{{{_DML}}}r'):
            rPr = run.find(f'{{{_DML}}}rPr')
            if rPr is not None and first_rPr is None:
                first_rPr = copy.deepcopy(rPr)
    for p in list(paras):
        txBody.remove(p)
    for bullet_text in (bullets or [""]):
        new_p = _etree.SubElement(txBody, f'{{{_DML}}}p')
        if first_pPr is not None:
            new_p.append(copy.deepcopy(first_pPr))
        new_r = _etree.SubElement(new_p, f'{{{_DML}}}r')
        if first_rPr is not None:
            new_r.append(copy.deepcopy(first_rPr))
        new_t = _etree.SubElement(new_r, f'{{{_DML}}}t')
        new_t.text = bullet_text
    return elem

def _find_best_base_slide_idx(prs) -> int:
    """
    Score each template slide and return the index of the one most suitable
    as a repeating base (title + body text, minimal complex groups/images).
    Prefer slides with ~2 text shapes, 0 text-bearing groups, 0 image refs.
    """
    best_idx, best_score = 0, float('inf')
    for i, slide in enumerate(prs.slides):
        cSld = slide._element.find(f'{{{_PML}}}cSld')
        if cSld is None:
            continue
        spTree = cSld.find(f'{{{_PML}}}spTree')
        if spTree is None:
            continue
        text_sp = 0
        text_groups = 0
        img_refs = 0
        for elem in spTree:
            tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
            if tag == 'sp' and _sp_has_text(elem) and not _is_slide_number_shape(elem):
                text_sp += 1
            elif tag == 'grpSp' and _elem_has_any_text(elem):
                text_groups += 1
            if _has_r_ref(elem):
                img_refs += 1
        score = abs(text_sp - 2) * 5 + text_groups * 15 + img_refs * 10
        if score < best_score:
            best_score = score
            best_idx = i
    return best_idx

def _is_freeform_sp(sp_elem) -> bool:
    """True if the shape is a freeform text box (Canva style — no <p:ph> element)."""
    if sp_elem is None:
        return False
    return sp_elem.find(f'.//{{{_PML}}}ph') is None

def _append_to_spTree(slide, elem) -> None:
    """Append an element to a slide's spTree (highest z-order = on top)."""
    cSld = slide._element.find(f'{{{_PML}}}cSld')
    if cSld is None:
        return
    spTree = cSld.find(f'{{{_PML}}}spTree')
    if spTree is not None:
        spTree.append(elem)

def _remove_all_sp_from_slide(slide) -> None:
    """Remove every <p:sp> from the slide's spTree (clears all placeholders and text boxes)."""
    cSld = slide._element.find(f'{{{_PML}}}cSld')
    if cSld is None:
        return
    spTree = cSld.find(f'{{{_PML}}}spTree')
    if spTree is None:
        return
    for elem in list(spTree):
        if elem.tag.split('}')[-1] == 'sp':
            spTree.remove(elem)

def _find_blank_layout_index(prs) -> int:
    """Return the index of the layout with the fewest placeholders (best for Canva-style)."""
    best, best_count = 0, float('inf')
    for i, layout in enumerate(prs.slide_layouts):
        try:
            count = sum(1 for _ in layout.placeholders)
            if count < best_count:
                best_count = count
                best = i
                if count == 0:
                    break
        except Exception:
            continue
    return best

# ---------- PPTX build (no overlap images) ----------
def build_presentation_from_plan(
    template_bytes: Optional[bytes],
    plan: Dict[str, Any],
    exact_reuse_images: bool = False,
) -> bytes:
    """
    Build a new deck:
      - If template provided, inherit its masters/layouts/styles.
      - If not, start from a blank Presentation().
      - If exact_reuse_images=True, copy PICTURE shapes from each template slide
        to the corresponding generated slide, but NEVER cover text:
          • add images BEFORE text (text stays on top)
          • auto-reposition/scale to a safe area if they overlap any text zone
    """
    prs = Presentation(io.BytesIO(template_bytes)) if template_bytes else Presentation()

    # Collect per-slide picture specs BEFORE clearing slides (if requested)
    template_pictures: List[List[Dict[str, any]]] = []
    if template_bytes and exact_reuse_images:
        for s in prs.slides:
            slide_specs: List[Dict[str, any]] = []
            for shape in s.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        slide_specs.append({
                            "blob": shape.image.blob,
                            "left": int(shape.left), "top": int(shape.top),
                            "width": int(shape.width), "height": int(shape.height),
                        })
                    except Exception:
                        pass
            template_pictures.append(slide_specs)

    # Collect per-slide <p:bg> elements BEFORE clearing (preserves colors/fills)
    template_bgs: List[Any] = []
    if template_bytes:
        for s in prs.slides:
            template_bgs.append(_get_slide_bg_elem(s))
    default_bg = next((bg for bg in template_bgs if bg is not None), None)

    # Collect decorative shapes, text shape positions, and layouts from each template slide
    template_decos: List[List[Any]] = []
    template_title_sps: List[Optional[Any]] = []
    template_body_sps: List[Optional[Any]] = []
    template_slide_layouts: List[Any] = []
    if template_bytes:
        for s in prs.slides:
            template_decos.append(_get_decorative_elems(s))
            t_sp, b_sp = _get_template_text_shapes(s)
            template_title_sps.append(t_sp)
            template_body_sps.append(b_sp)
            template_slide_layouts.append(s.slide_layout)

    # Use the cleanest template slide as the default fallback for all generated slides
    if template_bytes and prs.slides:
        base_idx = _find_best_base_slide_idx(prs)
    else:
        base_idx = 0
    default_deco      = template_decos[base_idx]     if template_decos     else []
    default_title_sp  = template_title_sps[base_idx] if template_title_sps else None
    default_body_sp   = template_body_sps[base_idx]  if template_body_sps  else None

    # Clear existing slides SAFELY (drop relationships first) so Office doesn't "repair"
    _clear_all_slides_safely(prs)

    # Fallback layout index for when there's no template or no per-slide layout available
    fallback_layout_idx = _find_title_and_content_layout_index(prs)
    if fallback_layout_idx is None:
        fallback_layout_idx = 1 if len(prs.slide_layouts) > 1 else 0

    slides_data = plan.get("slides", [])
    if not isinstance(slides_data, list) or not slides_data:
        raise RuntimeError("No slides found in plan.")

    slide_w, slide_h = int(prs.slide_width), int(prs.slide_height)

    for idx, s in enumerate(slides_data):
        title_txt = str(s.get("title", "")).strip()[:120] or f"Slide {idx+1}"
        bullets = [str(b).strip() for b in (s.get("bullets") or []) if str(b).strip()][:10]

        # Use the per-slide layout from the template (cycling) so each slide
        # inherits its unique background and visual style
        ti = (idx % max(1, len(template_slide_layouts))) if template_slide_layouts else 0
        if template_slide_layouts and ti < len(template_slide_layouts):
            slide_layout = template_slide_layouts[ti]
        else:
            slide_layout = prs.slide_layouts[fallback_layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # Apply template background (colour/gradient/fill) to this slide
        if template_bytes:
            bg = template_bgs[idx] if idx < len(template_bgs) else default_bg
            if bg is None:
                bg = default_bg
            if bg is not None:
                _apply_slide_bg(slide, bg)

        # Per-slide Canva shape references (cycle through template slides)
        if template_bytes:
            _d  = template_decos[ti]     if ti < len(template_decos)     else None
            _ts = template_title_sps[ti] if ti < len(template_title_sps) else None
            _bs = template_body_sps[ti]  if ti < len(template_body_sps)  else None
            deco     = _d  if _d  is not None else default_deco
            title_sp = _ts if _ts is not None else default_title_sp
            body_sp  = _bs if _bs is not None else default_body_sp
        else:
            deco = []; title_sp = None; body_sp = None

        # Insert decorative shapes behind everything
        if deco:
            _insert_decorative_elems(slide, deco)

        # --- Find placeholder rects early (for safe-zone selection)
        title_rect = None
        body_rect = None
        if slide.shapes.title:
            t = slide.shapes.title
            title_rect = _rect(int(t.left), int(t.top), int(t.width), int(t.height))
        for ph in slide.placeholders:
            try:
                if ph.is_placeholder and ph.placeholder_format.type not in (1,):  # 1 = title
                    body_rect = _rect(int(ph.left), int(ph.top), int(ph.width), int(ph.height))
                    break
            except Exception:
                continue

        # Gather ALL text-bearing zones (titles, subtitles, content placeholders, text boxes, etc.)
        text_zones = _collect_text_zones(slide)

        # --- 1) Insert images FIRST so text is on top (z-order safe) and avoid covering text
        if template_pictures and idx < len(template_pictures):
            for pic in template_pictures[idx]:
                target = _rect(pic["left"], pic["top"], pic["width"], pic["height"])
                if _overlaps_any_text(target, text_zones, thresh=0.10):
                    safe = _choose_safe_zone(
                        slide_w=slide_w, slide_h=slide_h,
                        title_rect=title_rect, body_rect=body_rect,
                        pad=91440  # ~0.1 inch in EMU
                    )
                    target = _fit_into_box(target, safe)
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(pic["blob"]),
                        target["left"], target["top"],
                        width=target["width"], height=target["height"]
                    )
                except Exception:
                    pass

        # --- 2) Add text
        _mg = 457200                        # 0.5 inch in EMU
        _tw = slide_w - 2 * _mg

        # Freeform = Canva-style shapes with no <p:ph> element.
        # Only use the clone path for freeform shapes — placeholder-based templates
        # (Google Slides, standard PPTX) must use the layout placeholder path so that
        # the layout's colors, fonts, and positions are applied correctly.
        is_freeform = _is_freeform_sp(title_sp) or _is_freeform_sp(body_sp)

        if is_freeform:
            # Canva-style: wipe layout placeholders and stamp our text into
            # clones of the template's own freeform text shapes.
            _remove_all_sp_from_slide(slide)

            if title_sp is not None:
                _append_to_spTree(slide, _clone_sp_with_text(title_sp, title_txt))
            else:
                tx = slide.shapes.add_textbox(_mg, int(slide_h * 0.05), _tw, int(slide_h * 0.17))
                tf = tx.text_frame
                tf.word_wrap = True
                tf.text = title_txt
                tf.paragraphs[0].font.size = Pt(28)
                tf.paragraphs[0].font.bold = True

            if body_sp is not None:
                _append_to_spTree(slide, _clone_sp_with_bullets(body_sp, bullets))
            else:
                body_top = int(slide_h * 0.25)
                tx = slide.shapes.add_textbox(_mg, body_top, _tw, slide_h - body_top - _mg)
                tf = tx.text_frame
                tf.word_wrap = True
                for i, bt in enumerate(bullets or [""]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = bt
                    p.font.size = Pt(18)

        else:
            # Standard placeholder path: works for Google Slides, standard PPTX,
            # and the no-template case. The layout provides correct styling.
            if slide.shapes.title:
                slide.shapes.title.text = title_txt
            else:
                tx = slide.shapes.add_textbox(_mg, int(slide_h * 0.05), _tw, int(slide_h * 0.17))
                tf = tx.text_frame
                tf.word_wrap = True
                tf.text = title_txt
                tf.paragraphs[0].font.size = Pt(28)
                tf.paragraphs[0].font.bold = True

            # Find best body placeholder:
            # - Exclude TITLE (1) and CENTER_TITLE (3) — they are title variants
            # - Prefer BODY (2) > SUBTITLE (4) > anything else
            #   so we never overwrite the title and always pick the main content area
            _TITLE_PH_TYPES = {1, 3}
            body_ph = None
            for preferred in (2, 4, None):
                for ph in slide.placeholders:
                    try:
                        ph_type = ph.placeholder_format.type
                        if ph_type in _TITLE_PH_TYPES:
                            continue
                        if preferred is None or ph_type == preferred:
                            body_ph = ph
                            break
                    except Exception:
                        continue
                if body_ph is not None:
                    break
            if body_ph:
                tf = body_ph.text_frame
                tf.clear()
                if bullets:
                    tf.paragraphs[0].text = bullets[0]
                    tf.paragraphs[0].level = 0
                    for bullet in bullets[1:]:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
                else:
                    tf.paragraphs[0].text = ""
            else:
                # Place textbox below the actual title placeholder to avoid overlap
                body_top = int(slide_h * 0.25)
                if slide.shapes.title:
                    try:
                        t = slide.shapes.title
                        body_top = int(t.top) + int(t.height) + _mg
                    except Exception:
                        pass
                body_top = min(body_top, int(slide_h * 0.75))  # never push off screen
                tx = slide.shapes.add_textbox(_mg, body_top, _tw, slide_h - body_top - _mg)
                tf = tx.text_frame
                tf.word_wrap = True
                for i, bt in enumerate(bullets or [""]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = bt
                    p.font.size = Pt(18)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

def _clear_all_slides_safely(prs: Presentation) -> None:
    """Delete all slides and drop relationships to prevent 'repair' prompts in Office."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.rId
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

def _find_title_and_content_layout_index(prs: Presentation) -> Optional[int]:
    """Heuristic: find a layout that has both a title and a body/content placeholder."""
    for i, layout in enumerate(prs.slide_layouts):
        has_title, has_body = False, False
        try:
            for ph in layout.placeholders:
                t = ph.placeholder_format.type
                if t == 1:      # Title
                    has_title = True
                if t in (2, 7): # Body or Content
                    has_body = True
            if has_title and has_body:
                return i
        except Exception:
            continue
    return None

# ---------- Geometry & text-zone helpers ----------
def _collect_text_zones(slide) -> List[Dict[str, int]]:
    """Rects for any shape likely to hold text: title/body/center-title/subtitle/content placeholders,
    plus any autoshape/textbox with a text_frame."""
    zones: List[Dict[str, int]] = []
    for sh in slide.shapes:
        try:
            if getattr(sh, "is_placeholder", False):
                ph_type = getattr(getattr(sh, "placeholder_format", None), "type", None)
                # common text placeholders: 1=TITLE, 2=BODY, 3=CENTER_TITLE, 4=SUBTITLE, 7=CONTENT
                if ph_type in (1, 2, 3, 4, 7):
                    zones.append(_rect(int(sh.left), int(sh.top), int(sh.width), int(sh.height)))
                    continue
            if getattr(sh, "has_text_frame", False):
                zones.append(_rect(int(sh.left), int(sh.top), int(sh.width), int(sh.height)))
        except Exception:
            continue
    return zones

def _rect(left: int, top: int, width: int, height: int) -> Dict[str, int]:
    return {"left": max(0, left), "top": max(0, top), "width": max(0, width), "height": max(0, height)}

def _intersect_area(a: Dict[str, int], b: Dict[str, int]) -> int:
    if not a or not b:
        return 0
    ax1, ay1, ax2, ay2 = a["left"], a["top"], a["left"] + a["width"], a["top"] + a["height"]
    bx1, by1, bx2, by2 = b["left"], b["top"], b["left"] + b["width"], b["top"] + b["height"]
    ix1, iy1, ix2, iy2 = max(ax1, bx1), max(ay1, by1), min(ax2, bx2), min(ay2, by2)
    if ix2 <= ix1 or iy2 <= iy1:
        return 0
    return (ix2 - ix1) * (iy2 - iy1)

def _overlaps_any_text(img: Dict[str, int], zones: List[Dict[str, int]], thresh: float = 0.10) -> bool:
    area = max(1, img["width"] * img["height"])
    for z in zones:
        if _intersect_area(img, z) / area > thresh:
            return True
    return False

def _choose_safe_zone(slide_w: int, slide_h: int,
                      title_rect: Optional[Dict[str, int]],
                      body_rect: Optional[Dict[str, int]],
                      pad: int = 0) -> Dict[str, int]:
    """
    Prefer a column to the RIGHT of the body; if too narrow, use BELOW the body.
    If no body placeholder, fall back to area under the title; else a right sidebar.
    """
    if body_rect:
        # Right of body
        right_left = body_rect["left"] + body_rect["width"] + pad
        right_width = max(0, slide_w - right_left - pad)
        right_top = body_rect["top"]
        right_height = body_rect["height"]
        if right_width >= slide_w * 0.18 and right_height >= slide_h * 0.18:
            return _rect(right_left, right_top, right_width, right_height)

        # Below body
        below_top = body_rect["top"] + body_rect["height"] + pad
        below_height = max(0, slide_h - below_top - pad)
        if below_height >= slide_h * 0.18:
            return _rect(pad, below_top, max(0, slide_w - 2 * pad), below_height)

        # Left of body (last resort)
        left_width = max(0, body_rect["left"] - 2 * pad)
        if left_width >= slide_w * 0.18:
            return _rect(pad, body_rect["top"], left_width, body_rect["height"])

    if title_rect:
        area_top = title_rect["top"] + title_rect["height"] + pad
        area_height = max(0, slide_h - area_top - pad)
        return _rect(pad, area_top, max(0, slide_w - 2 * pad), area_height)

    # Fallback: right sidebar
    sidebar_left = int(slide_w * 0.64) + pad
    sidebar_width = max(0, int(slide_w * 0.36) - 2 * pad)
    sidebar_top = int(slide_h * 0.18) + pad
    sidebar_height = max(0, int(slide_h * 0.72) - 2 * pad)
    return _rect(sidebar_left, sidebar_top, sidebar_width, sidebar_height)

def _fit_into_box(img: Dict[str, int], box: Dict[str, int]) -> Dict[str, int]:
    """Scale img to fit within box, keep aspect ratio; center it inside the box."""
    iw, ih = max(1, img["width"]), max(1, img["height"])
    bw, bh = max(1, box["width"]), max(1, box["height"])
    scale = min(bw / iw, bh / ih, 1.0)
    nw, nh = int(iw * scale), int(ih * scale)
    nl = box["left"] + (bw - nw) // 2
    nt = box["top"] + (bh - nh) // 2
    return _rect(nl, nt, nw, nh)

